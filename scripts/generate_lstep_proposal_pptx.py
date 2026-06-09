#!/usr/bin/env python3
"""Generate L-step premium proposal PowerPoint for SANN."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Colors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xFA)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
GREEN = RGBColor(0x06, 0xC7, 0x55)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Hiragino Sans"
TOTAL_SLIDES = 12
OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "L-step-導入詳細提案資料_SANN.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb_hex(c: RGBColor) -> str:
    return f"{c}"


def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, page: int):
    left = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(6), Inches(0.35))
    p = left.text_frame.paragraphs[0]
    p.text = "株式会社SANN 御中"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = FONT

    right = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.2), Inches(0.35))
    p = right.text_frame.paragraphs[0]
    p.text = f"{page} / {TOTAL_SLIDES}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_LIGHT
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def add_title_bar(slide, title, section=None):
    if section:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(1.2), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = section
        p.font.size = Pt(36)
        p.font.color.rgb = GRAY_LIGHT
        p.font.name = FONT
        p.font.bold = True

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = DARK
    p.font.name = FONT
    p.font.bold = True

    line = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.65), Inches(1.5), Inches(0.03)
    )  # rectangle
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def add_notes(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_text_box(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size=14,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(6)
    return box


def add_card(slide, left, top, width, height, number, title, body_lines):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    num_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width, Inches(0.4))
    p = num_box.text_frame.paragraphs[0]
    p.text = number
    p.font.size = Pt(22)
    p.font.color.rgb = GREEN
    p.font.name = FONT
    p.font.bold = True

    title_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), Inches(0.5)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.font.name = FONT
    p.font.bold = True

    body_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(1.05), width - Inches(0.4), height - Inches(1.2)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.name = FONT


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_LIGHT
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = DARK
            p.font.name = FONT

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.font.name = FONT
                if j == len(row) - 1 and i == len(rows) - 1:
                    p.font.bold = True
                    p.font.color.rgb = GREEN
    return table


def slide_01(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide,
        Inches(0.6),
        Inches(1.8),
        Inches(12),
        Inches(0.5),
        ["PROPOSAL"],
        size=12,
        color=GRAY,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(2.5),
        Inches(12),
        Inches(1.2),
        ["Lステップ導入詳細ご説明資料"],
        size=36,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(3.8),
        Inches(12),
        Inches(1),
        [
            "プレミアムプラン導入に向けた",
            "費用・設計スケジュールのご説明",
        ],
        size=22,
        color=GRAY,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(6.2),
        Inches(6),
        Inches(0.6),
        ["株式会社SANN 御中", "2026年5月"],
        size=14,
        color=GRAY,
    )
    add_footer(slide, 1)
    add_notes(
        slide,
        "前回は全体構想をお話ししました。本日はプレミアムプランについて、"
        "費用の内訳と要件定義完了を前提とした導入スケジュールを具体的にお伝えします。",
    )


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "本日のご説明内容")
    add_card(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(3.9),
        Inches(3.8),
        "01",
        "導入費用の詳細確認",
        [
            "初期構築費（プレミアム",
            "¥1,400,000）の",
            "項目ごとの内訳をご説明",
        ],
    )
    add_card(
        slide,
        Inches(4.7),
        Inches(2.0),
        Inches(3.9),
        Inches(3.8),
        "02",
        "導入スケジュールの確認",
        [
            "要件定義完了を前提に",
            "設計・制作4週間＋",
            "テスト・運用開始2週間",
            "計6週間で進行",
        ],
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(2.0),
        Inches(3.9),
        Inches(3.8),
        "03",
        "運用開始までの流れ確認",
        [
            "御社側のご準備事項と",
            "当社サポート範囲を整理",
        ],
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(6.1),
        Inches(12),
        Inches(0.6),
        [
            "前回ご説明した全体構想を前提に、"
            "「導入後のイメージ」と「投資判断材料」を具体化します"
        ],
        size=13,
        color=GRAY,
    )
    add_footer(slide, 2)
    add_notes(
        slide,
        "要件定義は前回までで完了しているため、本日は残り6週間で何をするかと"
        "140万円の内訳に焦点を当てます。売り込みではなく、導入後の見通しを固める場にしたいです。",
    )


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "プレミアムプランで実現できること")

    add_text_box(
        slide,
        Inches(0.6),
        Inches(2.0),
        Inches(5.5),
        Inches(3.5),
        [
            "✓  応募者教育（自動配信シナリオ）",
            "✓  説明会自動化（日程調整・リマインド）",
            "✓  属性別シナリオ（学部・志望・ステータス別）",
            "✓  診断機能（適性・興味の可視化）",
            "✓  採用資産化（友だち・行動データの蓄積）",
        ],
        size=14,
    )

    headers = ["", "スタンダード", "プレミアム"]
    rows = [
        ["応募者教育", "○", "◎"],
        ["説明会自動化", "○", "◎"],
        ["属性別シナリオ", "○", "◎"],
        ["診断機能", "△", "◎"],
        ["分析・改善", "○", "◎"],
    ]
    add_table(
        slide,
        Inches(6.5),
        Inches(2.0),
        Inches(6.2),
        Inches(2.8),
        headers,
        rows,
        [Inches(2.8), Inches(1.7), Inches(1.7)],
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(5.8),
        Inches(12),
        Inches(0.8),
        [
            "プレミアム＝「採用の一時対応」ではなく「採用資産を積み上げる仕組み」",
        ],
        size=16,
        color=GREEN,
        bold=True,
    )
    add_footer(slide, 3)
    add_notes(
        slide,
        "説明会自動化はスタンダードでも実現できます。プレミアムは属性別シナリオの高度化、"
        "本格的な診断機能、分析改善の伴走まで含む構成です。",
    )


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "初期構築費の内訳（プレミアムプラン）")

    headers = ["項目", "内容概要", "金額（税別）"]
    rows = [
        ["戦略設計", "要件定義済み内容の反映・全体設計", "¥100,000"],
        ["シナリオ構築", "教育・nurture・選考連携の構築", "¥220,000"],
        ["L-CAST構築", "リッチメニュー・動画・フォーム基盤", "¥50,000"],
        ["診断機能", "質問設計・分岐・結果ページの実装", "¥120,000"],
        ["外部連携", "フォーム・CRM・スプレッドシート等", "¥300,000"],
        ["クリエイティブ作成", "リッチメニュー・Lフレックス・診断素材等", "¥260,000"],
        ["タグ設計", "属性・行動・ステータス管理", "¥80,000"],
        ["分析環境構築", "GA4 / レポート環境の整備", "¥120,000"],
        ["テスト・調整・納品", "動作確認・修正・運用引き継ぎ", "¥150,000"],
        ["合計", "", "¥1,400,000"],
    ]
    add_table(
        slide,
        Inches(0.6),
        Inches(1.9),
        Inches(12.1),
        Inches(4.5),
        headers,
        rows,
        [Inches(2.0), Inches(6.3), Inches(1.8)],
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(6.5),
        Inches(12),
        Inches(0.5),
        [
            "※ お支払い：契約時50%（¥700,000）／運用開始時50%（¥700,000）",
            "※ 要件定義は完了済みのため、戦略設計は反映・調整中心の工数です",
        ],
        size=10,
        color=GRAY,
    )
    add_footer(slide, 4)
    add_notes(
        slide,
        "140万円のうち、外部連携30万円とクリエイティブ26万円が制作実務の中心です。"
        "シナリオ22万円・診断12万円・L-CAST5万円は、前回固めた要件をもとに構築します。",
    )


def slide_05(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "月額費用の内訳")

    headers = ["項目", "サービス内容", "月額（税別）"]
    rows = [
        ["公式LINE", "公式アカウント・配信基盤", "¥5,000"],
        ["Lステップ", "プレミアムプラン利用料", "¥32,780"],
        ["L-CAST", "リッチメニュー・動画配信", "¥14,960"],
        ["運用保守", "月次改善・障害対応・相談", "¥100,000"],
        ["月額合計（ツールのみ）", "", "約 ¥52,740"],
        ["月額合計（ツール＋運用保守）", "", "約 ¥152,740"],
    ]
    add_table(
        slide,
        Inches(0.6),
        Inches(1.9),
        Inches(8.5),
        Inches(3.2),
        headers,
        rows,
        [Inches(1.8), Inches(4.2), Inches(2.5)],
    )

    highlight = slide.shapes.add_shape(1, Inches(9.4), Inches(2.2), Inches(3.3), Inches(2.8))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = BG_LIGHT
    highlight.line.color.rgb = GREEN
    highlight.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(9.6),
        Inches(2.4),
        Inches(3),
        Inches(2.4),
        [
            "初年度イメージ（税別）",
            "",
            "初期構築",
            "¥1,400,000",
            "",
            "＋ 月額 × 12ヶ月",
            "",
            "＝ 約 ¥3,232,880",
        ],
        size=13,
        bold=False,
    )
    tf = slide.shapes[-1].text_frame
    for i, p in enumerate(tf.paragraphs):
        if "¥3,232,880" in p.text or "¥1,400,000" in p.text:
            p.font.bold = True
            p.font.color.rgb = GREEN
            p.font.size = Pt(16 if "3,232" in p.text else 14)

    add_text_box(
        slide,
        Inches(0.6),
        Inches(6.3),
        Inches(12),
        Inches(0.4),
        ["※ 運用保守は内製化後に縮小可能です"],
        size=10,
        color=GRAY,
    )
    add_footer(slide, 5)
    add_notes(
        slide,
        "月額はツールのみ約5.3万円、運用サポート込みで約15.3万円です。"
        "初年度は初期140万円を含め約323万円が目安です。",
    )


def slide_06(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "投資対効果の考え方")

    # Bar comparison (simplified)
    bar1 = slide.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.45))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = GRAY_LIGHT
    bar1.line.fill.background()
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.95),
        Inches(6),
        Inches(0.35),
        ["採用エージェント経由　1名 ¥500,000〜¥1,000,000"],
        size=12,
        color=GRAY,
    )

    bar2 = slide.shapes.add_shape(1, Inches(0.8), Inches(3.2), Inches(4.2), Inches(0.45))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = GREEN
    bar2.line.fill.background()
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.85),
        Inches(6),
        Inches(0.35),
        ["LINE経由（プレミアム）初年度 約¥3,230,000"],
        size=12,
        color=GRAY,
    )
    add_text_box(
        slide,
        Inches(0.8),
        Inches(3.75),
        Inches(6),
        Inches(0.35),
        ["2年目以降 約¥1,830,000/年（初期費用なし）"],
        size=11,
        color=GRAY,
    )

    add_text_box(
        slide,
        Inches(7.0),
        Inches(2.0),
        Inches(5.5),
        Inches(3.5),
        [
            "① 採用単価の削減",
            "　エージェント依存の低減",
            "",
            "② 工数削減",
            "　説明会・個別案内の自動化",
            "",
            "③ 中長期の資産化",
            "　友だち・行動・診断データの蓄積",
        ],
        size=14,
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(5.5),
        Inches(12),
        Inches(1),
        [
            "損益分岐の目安：エージェント経由 3〜4名分の採用コスト削減で",
            "初年度投資を回収可能（御社の採用規模により変動）",
        ],
        size=13,
        color=GRAY,
    )
    add_footer(slide, 6)
    add_notes(
        slide,
        "初年度は初期費用が140万円のため投資額は大きくなりますが、"
        "診断・外部連携・クリエイティブまで一括で完成度を上げられます。",
    )


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "導入スケジュール（残り 約6週間）")

    # Completed bar
    done = slide.shapes.add_shape(1, Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.55))
    done.fill.solid()
    done.fill.fore_color.rgb = BG_LIGHT
    done.line.color.rgb = GRAY_LIGHT
    add_text_box(
        slide,
        Inches(0.8),
        Inches(2.08),
        Inches(10),
        Inches(0.4),
        ["✓  要件定義 … 完了（前回までに実施済み）"],
        size=14,
        color=GRAY_LIGHT,
    )

    phases = [
        ("Week 1-2", "設計", "Phase 2"),
        ("Week 3-4", "制作", "Phase 3"),
        ("Week 5", "テスト", "Phase 4"),
        ("Week 6", "運用開始", "Phase 5"),
    ]
    x_start = 0.8
    for i, (week, name, phase) in enumerate(phases):
        x = Inches(x_start + i * 3.1)
        circle = slide.shapes.add_shape(9, x, Inches(3.2), Inches(0.35), Inches(0.35))  # oval
        circle.fill.solid()
        circle.fill.fore_color.rgb = GREEN
        circle.line.fill.background()

        line = slide.shapes.add_shape(1, x + Inches(0.15), Inches(3.0), Inches(2.8), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = BORDER if i < 3 else BORDER
        line.line.fill.background()

        add_text_box(slide, x - Inches(0.1), Inches(2.75), Inches(2.8), Inches(0.35), [week], size=11, color=GRAY)
        add_text_box(slide, x - Inches(0.1), Inches(3.65), Inches(2.8), Inches(0.4), [name], size=16, bold=True)
        add_text_box(slide, x - Inches(0.1), Inches(4.05), Inches(2.8), Inches(0.35), [phase], size=10, color=GRAY)

    add_text_box(
        slide,
        Inches(10.5),
        Inches(2.5),
        Inches(2.2),
        Inches(1),
        ["6", "週間"],
        size=48,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(5.0),
        Inches(12),
        Inches(1.2),
        [
            "★ Week 2末：設計書承認",
            "★ Week 4末：制作完了・受入テスト開始",
            "★ Week 6：運用開始",
            "",
            "※ 要件定義完了により、従来10週間想定を6週間に短縮",
        ],
        size=13,
    )
    add_footer(slide, 7)
    add_notes(
        slide,
        "前回までで要件定義は完了しているため、本日から約6週間で運用開始まで進められます。",
    )


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Phase 2｜設計（Week 1〜2）")

    ribbon = slide.shapes.add_shape(1, Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.4))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = BG_LIGHT
    add_text_box(
        slide,
        Inches(0.8),
        Inches(1.78),
        Inches(10),
        Inches(0.35),
        ["前提：要件定義は完了済み。定義書をもとに設計・制作へ移行します。"],
        size=11,
        color=GRAY_LIGHT,
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(2.3),
        Inches(5.8),
        Inches(3.8),
        [
            "【当社の作業】",
            "・シナリオ設計書・画面遷移図の作成",
            "・属性別シナリオ（分岐・タグ連動）の設計",
            "・L-CAST設計（リッチメニュー・動画・フォーム）",
            "・診断機能の質問・分岐・結果画面の設計",
            "・外部連携（フォーム・CRM等）の接続設計",
            "・分析設計（計測項目・レポート形式）",
        ],
        size=12,
    )

    box = slide.shapes.add_shape(1, Inches(6.8), Inches(2.3), Inches(5.9), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xFD, 0xF4)
    box.line.color.rgb = GREEN
    box.line.width = Pt(1)

    add_text_box(
        slide,
        Inches(7.0),
        Inches(2.45),
        Inches(5.5),
        Inches(3.2),
        [
            "【御社のご準備】",
            "・教育コンテンツ素材（動画・PDF等）の共有",
            "・説明会日程・会場・定員情報の整理",
            "・クリエイティブの方向性確認",
            "・設計書へのフィードバック（48h以内推奨）",
            "",
            "【成果物】",
            "設計書一式（シナリオ/L-CAST/診断/連携/分析）",
        ],
        size=12,
    )

    add_text_box(
        slide,
        Inches(0.6),
        Inches(6.0),
        Inches(12),
        Inches(0.4),
        ["★ Week 2末：設計書承認 → Week 3から制作開始"],
        size=13,
        color=GREEN,
        bold=True,
    )
    add_footer(slide, 8)
    add_notes(slide, "設計2週間は、前回の要件定義を実装可能な設計書に落とし込む期間です。")


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "Phase 3〜5｜制作 〜 運用開始（Week 3〜6）")

    blocks = [
        (
            "Phase 3｜制作（Week 3〜4）",
            [
                "・シナリオ構築（教育・説明会・選考連携）",
                "・L-CAST実装・診断機能の実装",
                "・外部連携の設定・接続テスト",
                "・クリエイティブ制作・組み込み",
                "・タグ設計・分析環境の構築",
                "★ Week 4末：制作完了",
            ],
        ),
        (
            "Phase 4｜テスト（Week 5）",
            [
                "・結合テスト（登録→教育→説明会→選考）",
                "・御社担当者による受入テスト",
                "・修正・最終調整",
                "・社内レクチャー（2時間）",
            ],
        ),
        (
            "Phase 5｜運用開始（Week 6〜）",
            [
                "・本番環境への切替",
                "・初回配信・説明会募集の開始",
                "・運用マニュアル納品",
                "・初月の伴走サポート",
                "★ Week 6：運用開始",
            ],
        ),
    ]

    y = 2.0
    for i, (title, lines) in enumerate(blocks):
        left_border = slide.shapes.add_shape(
            1, Inches(0.55), Inches(y), Inches(0.08), Inches(1.35 if i < 2 else 1.2)
        )
        left_border.fill.solid()
        left_border.fill.fore_color.rgb = GREEN if i >= 1 else GRAY_LIGHT
        left_border.line.fill.background()

        add_text_box(
            slide,
            Inches(0.75),
            Inches(y),
            Inches(11.5),
            Inches(1.5),
            [title, *lines],
            size=12 if i > 0 else 12,
        )
        y += 1.55

    add_footer(slide, 9)
    add_notes(slide, "制作は集中的に2週間。Week 5はテスト、Week 6で本番切替です。")


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "導入後の学生導線イメージ")

    nodes = ["SNS・HP", "LINE登録", "自動教育", "説明会", "選考"]
    x = 0.7
    for i, label in enumerate(nodes):
        nx = Inches(x + i * 2.45)
        node = slide.shapes.add_shape(1, nx, Inches(2.8), Inches(2.0), Inches(0.7))
        node.fill.solid()
        node.fill.fore_color.rgb = WHITE
        node.line.color.rgb = BORDER
        p = slide.shapes.add_textbox(nx, Inches(2.95), Inches(2.0), Inches(0.5)).text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(13)
        p.font.name = FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if i < len(nodes) - 1:
            arrow = slide.shapes.add_textbox(
                nx + Inches(2.05), Inches(2.95), Inches(0.35), Inches(0.4)
            )
            arrow.text_frame.paragraphs[0].text = "→"
            arrow.text_frame.paragraphs[0].font.size = Pt(18)
            arrow.text_frame.paragraphs[0].font.color.rgb = GREEN

    details = [
        "① 登録 … QR / SNS連携。学部・学年・志望をタグ取得",
        "② 教育 … 属性別シナリオで3〜5日間のステップ配信",
        "③ 説明会 … 希望日時選択、前日・当日リマインド自動送信",
        "④ 選考 … 参加者へ個別案内、未参加者へ再ナーチャリング",
    ]
    add_text_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(1.8), details, size=13, color=GRAY)

    add_text_box(
        slide,
        Inches(0.6),
        Inches(5.8),
        Inches(12),
        Inches(0.4),
        ["→ 全導線がLステップ上で可視化・改善可能"],
        size=15,
        color=GREEN,
        bold=True,
    )
    add_footer(slide, 10)
    add_notes(slide, "属性別シナリオにより、学部や志望度に応じた最適な教育・案内が自動で届きます。")


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title_bar(slide, "運用開始後のサポート")
    add_text_box(
        slide,
        Inches(0.6),
        Inches(1.72),
        Inches(12),
        Inches(0.35),
        ["月額運用保守 ¥100,000 に含まれる"],
        size=12,
        color=GRAY,
    )

    items = [
        ("月次改善MTG", "配信結果の振り返り\n改善施策の提案・実装"),
        ("データ分析", "開封率・遷移率\n説明会参加率・選考移行率"),
        ("配信・シナリオ改善", "A/Bテスト\nシナリオ追加・修正（月2回まで）"),
        ("相談対応", "チャット・メール随時\n障害の一次対応"),
    ]
    for i, (title, body) in enumerate(items):
        x = Inches(0.6 + i * 3.15)
        add_card(slide, x, Inches(2.1), Inches(2.95), Inches(3.5), f"0{i+1}", title, body.split("\n"))

    add_text_box(
        slide,
        Inches(0.6),
        Inches(5.9),
        Inches(12),
        Inches(0.6),
        [
            "「作って終わり」ではなく、",
            "採用成果が出るまで伴走するパートナー契約です",
        ],
        size=16,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 11)
    add_notes(slide, "プレミアムの分析・改善は、運用開始後も月次MTGで継続します。")


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(
        slide,
        Inches(0.6),
        Inches(1.5),
        Inches(12),
        Inches(0.6),
        ["まとめ"],
        size=32,
        bold=True,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(2.5),
        Inches(12),
        Inches(1.5),
        [
            "本提案は「採用活動の効率化」ではなく、",
            "「採用資産の構築」です。",
        ],
        size=24,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide,
        Inches(0.6),
        Inches(4.2),
        Inches(12),
        Inches(0.8),
        [
            "初期構築 ¥1,400,000（税別）",
            "設計・制作 4週間 ＋ テスト・運用開始 2週間 ＝ 約6週間で運用開始",
        ],
        size=16,
        color=GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide,
        Inches(1.5),
        Inches(5.2),
        Inches(10),
        Inches(1.5),
        [
            "次のステップ",
            "① ご質問・ご懸念点の確認",
            "② 導入開始時期のすり合わせ（例：今週キックオフ → 6週間後に運用開始）",
            "③ 契約・制作着手日程の調整",
            "",
            "担当：＿＿＿＿　／　お打ち合わせ調整：＿＿＿＿",
        ],
        size=14,
    )
    add_footer(slide, 12)
    add_notes(
        slide,
        "要件定義は完了しているので、あとは6週間で形にする段階です。"
        "御社の採用カレンダーに合わせ、キックオフ時期をご一緒に決めさせてください。",
    )


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_01,
        slide_02,
        slide_03,
        slide_04,
        slide_05,
        slide_06,
        slide_07,
        slide_08,
        slide_09,
        slide_10,
        slide_11,
        slide_12,
    ]
    for fn in builders:
        fn(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
